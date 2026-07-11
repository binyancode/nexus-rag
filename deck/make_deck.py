# -*- coding: utf-8 -*-
"""Assemble the 法规检索系统设计 deck: native editable text slides + diagram images
(with native editable titles). Run: python make_deck.py"""
import os
import subprocess

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

import diagrams as D
from slides_lib import html_wrap

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "assets")
os.makedirs(OUT, exist_ok=True)
PPTX = os.path.join(os.path.dirname(HERE), "法规检索系统设计.pptx")
CHROME = r"C:\Program Files\Google\Chrome\Application\chrome.exe"
if not os.path.exists(CHROME):
    CHROME = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"

PXE = 9525
def PX(v): return Emu(int(round(v * PXE)))
def PT(v): return Pt(v * 0.75)
def C(h): return RGBColor.from_string(h.lstrip("#"))
YAHEI = "Microsoft YaHei"

BG="#EEF3FA"; PANEL="#FFFFFF"; INK="#17222E"; MUTE="#5A6B7B"; FAINT="#93A1B3"
NAVY="#123A63"; BLUE="#1470C4"; BLUE_SOFT="#DBEAFA"; TEAL="#0E9C9C"; TEAL_SOFT="#D3EFEE"
GREEN="#2E9E6B"; GREEN_SOFT="#D7F0E3"; ORANGE="#E0812E"; ORANGE_SOFT="#FBE9D6"
PURPLE="#7A5CC0"; PURPLE_SOFT="#E7E0F6"; RED="#D2503F"; LINE_C="#DCE4EE"; BORDER="#E4EBF4"


def rrect(s, x, y, w, h, fill, line=None, rad=14, lw=1.0):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(x), PX(y), PX(w), PX(h))
    try:
        sp.adjustments[0] = max(0.0, min(0.5, rad / float(min(w, h))))
    except Exception:
        pass
    sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    if line:
        sp.line.color.rgb = C(line); sp.line.width = Pt(lw)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def rct(s, x, y, w, h, fill, line=None, lw=1.0):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PX(x), PX(y), PX(w), PX(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    if line:
        sp.line.color.rgb = C(line); sp.line.width = Pt(lw)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def hline(s, x1, y, x2, color, weight=1.2):
    return rct(s, x1, y, x2 - x1, weight, color)


def txt(s, x, y, w, h, t, size, color, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font=YAHEI):
    tb = s.shapes.add_textbox(PX(x), PX(y), PX(w), PX(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, ln in enumerate(str(t).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        f = r.font; f.size = PT(size); f.bold = bold; f.name = font; f.color.rgb = C(color)
    return tb


def circle(s, x, y, d, fill, line=None, lw=1.4):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, PX(x), PX(y), PX(d), PX(d))
    sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    if line:
        sp.line.color.rgb = C(line); sp.line.width = Pt(lw)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def bg_native(s):
    rct(s, 0, 0, 1280, 720, BG)


def header_native(s, kicker, title, idx):
    rct(s, 0, 0, 520, 8, BLUE); rct(s, 520, 0, 380, 8, TEAL); rct(s, 900, 0, 380, 8, GREEN)
    txt(s, 70, 40, 700, 24, kicker, 16, TEAL, True)
    txt(s, 70, 64, 1050, 50, title, 32, NAVY, True)
    hline(s, 70, 126, 1210, LINE_C, 2)


def footer_native(s, label="法规检索系统设计 · 两层图 + SQG 编译"):
    hline(s, 70, 684, 1210, LINE_C, 1)
    txt(s, 70, 690, 700, 20, label, 12.5, FAINT)
    txt(s, 510, 690, 700, 20, "v0.2 · 讨论稿", 12.5, FAINT, False, PP_ALIGN.RIGHT)


def card_native(s, x, y, w, h, accent):
    rrect(s, x, y, w, h, PANEL, line=BORDER, rad=14, lw=1.0)
    rrect(s, x, y, 6, h, accent, rad=3)


# ---------------- native slides ----------------
def s_philosophy(s, idx):
    bg_native(s); header_native(s, "PRINCIPLES", "核心设计理念", idx)
    txt(s, 70, 140, 1140, 28, "把检索从“黑箱猜”变成“看得见的图上计算”：大模型编译可见的查询步骤，在两层图上执行，条条可溯源。", 16, MUTE)
    cards = [
        ("01 可见的算子图(SQG)", "大模型把“怎么查”编译成算子 DAG，每步一句人话，看得见、可排查。", PURPLE),
        ("02 两层图", "实体层管结构关系（唯一概念节点），块层管原文与向量入口，evidence 边相连。", BLUE),
        ("03 逻辑与执行分离", "大模型只说“想干什么”；“怎么做”交给优化器/执行器，不碰底层零件。", TEAL),
        ("04 关系只用 LLM 建", "便宜手段只提名候选，判定权全给 LLM；关系存一次，反向靠遍历。", ORANGE),
        ("05 优化后再执行", "校验→绑定→下推/并行/惰性取原文/缓存，同一逻辑图跑得又快又省。", GREEN),
        ("06 可溯源", "每条答案带 fullname 出处；执行后每步回填实际输出，便于定位问题。", NAVY),
    ]
    x0, y0, cw, ch, gx, gy = 70, 182, 366, 152, 21, 20
    for i, (t, d, c) in enumerate(cards):
        col, row = i % 3, i // 3
        x, y = x0 + col * (cw + gx), y0 + row * (ch + gy)
        card_native(s, x, y, cw, ch, c)
        txt(s, x + 24, y + 24, cw - 40, 28, t, 17, INK, True)
        txt(s, x + 24, y + 60, cw - 44, 84, d, 13.8, MUTE)
    footer_native(s)


def s_logical_ops(s, idx):
    bg_native(s); header_native(s, "OPERATORS", "逻辑算子目录（大模型写的）", idx)
    txt(s, 70, 140, 1140, 28, "大模型只能从这一小把“意图级”算子里挑，一句人话能说清；算子间依赖=思考顺序。固定目录，不让 LLM 造新算子。", 16, MUTE)
    ops = [
        ("检索 Retrieve", "找出与某描述相关的内容/对象（如“IND 需要的法规”）", BLUE),
        ("关联 Relate", "从已有结果顺关系找相关对象（如“适用于哪些申报”“被谁替代”）", TEAL),
        ("筛选 Filter", "从结果里按条件留一部分（如“只留现行有效的”）", GREEN),
        ("对比 Compare", "比两组的异同 / 找各自独有", ORANGE),
        ("汇总 Summarize", "概括归纳一堆内容", PURPLE),
        ("校验 Verify", "检查结果全不全、对不对（可有界回环重来）", RED),
        ("回答 Answer", "综合成答案并标出处（fullname 溯源）", NAVY),
    ]
    y = 184
    for i, (t, d, c) in enumerate(ops):
        card_native(s, 70, y, 1140, 54, c)
        txt(s, 94, y, 240, 54, t, 17, INK, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        txt(s, 350, y, 840, 54, d, 14.5, MUTE, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        y += 62
    footer_native(s)


def s_physical_map(s, idx):
    bg_native(s); header_native(s, "MAPPING", "物理算子 + 逻辑→物理映射", idx)
    txt(s, 70, 140, 1140, 28, "优化器把每个逻辑算子展开成一串物理算子（执行零件）；一个逻辑算子→多个物理算子，且可因输入不同有不同展开。", 15.5, MUTE)
    # left: physical catalog
    card_native(s, 70, 180, 470, 452, BLUE)
    txt(s, 94, 198, 430, 26, "物理算子（执行零件，大模型不可见）", 16, INK, True)
    cat = [
        ("进门取数", "Resolve(精确)｜Seed(向量)｜ScanEntities"),
        ("图导航", "Lift｜Ground｜Traverse(带权扩散)｜StructUp/Down"),
        ("集合过滤", "Intersect｜Diff｜Union｜Dedup｜Filter"),
        ("排序裁剪", "Rerank｜TopK｜SortByFullname"),
        ("LLM 加工", "Summarize｜Extract｜Compare｜Generate｜Judge"),
        ("校验控制", "SetCheck｜ProvenanceCheck｜Threshold｜Replan"),
    ]
    y = 236
    for t, d in cat:
        txt(s, 94, y, 130, 22, "· " + t, 14, BLUE, True)
        txt(s, 210, y, 320, 40, d, 12.8, MUTE, font="Consolas")
        y += 66
    # right: mapping
    card_native(s, 560, 180, 650, 452, ORANGE)
    txt(s, 584, 198, 600, 26, "逻辑算子 → 物理子计划", 16, INK, True)
    mp = [
        ("检索", "(精确) Resolve→Traverse→Dedup  ｜ (语义) Seed→Lift→Dedup"),
        ("关联", "[Lift]→Traverse(带权扩散)→[Ground]→Dedup"),
        ("筛选", "FilterField(结构化) ｜ FilterSemantic(相似度/LLM)"),
        ("对比", "Intersect｜Diff｜Union  (＋Ground→LLM_Compare)"),
        ("汇总", "Ground→[分组]→LLM_Summarize"),
        ("校验", "SetCheck ＋ ProvenanceCheck ＋ Threshold(→Replan)"),
        ("回答", "Ground→SortByFullname→Dedup→LLM_Generate(cite)"),
    ]
    y = 238
    for t, d in mp:
        rrect(s, 584, y, 46, 30, ORANGE_SOFT, rad=7)
        txt(s, 584, y, 46, 30, t, 13, ORANGE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        txt(s, 644, y, 550, 34, d, 12.6, MUTE, font="Consolas")
        y += 54
    footer_native(s)


def s_optimizer(s, idx):
    bg_native(s); header_native(s, "OPTIMIZER", "优化器做什么", idx)
    txt(s, 70, 140, 1140, 28, "优化器把大模型的逻辑图，编译成又对又快的可执行 DAG。", 16, MUTE)
    steps = [
        ("① 校验 / 纠错", "算子输入输出合法？必须无环？参数合规？非法就打回或自动修。", BLUE),
        ("② 绑定物理实现", "把每个逻辑算子展开成一串执行零件（进门二选一：Resolve / Seed）。", TEAL),
        ("③ 优化", "谓词下推、公共子表达式复用、独立分支并行、算子融合、惰性取原文、缓存。", ORANGE),
        ("④ 编排执行", "拓扑序执行；校验不过可有界回环重规划。", GREEN),
    ]
    y = 186
    for t, d, c in steps:
        card_native(s, 70, y, 1140, 74, c)
        txt(s, 94, y + 14, 300, 26, t, 17, INK, True)
        txt(s, 94, y + 44, 1090, 24, d, 14, MUTE)
        y += 88
    rrect(s, 70, y + 4, 1140, 60, "#0E2E52", rad=12)
    rrect(s, 70, y + 4, 6, 60, ORANGE, rad=3)
    txt(s, 96, y + 14, 1090, 44, "惰性取原文＝中间只用 id 求集合，最终幸存者才去块存储取原文；同一张逻辑图，因优化跑得更快更省。",
        14.5, "#DCEBFA", True)
    footer_native(s)


def s_compare(s, idx):
    bg_native(s); header_native(s, "SUMMARY", "两个例子对比", idx)
    txt(s, 70, 140, 1140, 28, "同一套算子与执行器，按“入口是精确名还是语义描述”自动选择是否查向量。", 16, MUTE)
    cols = ["", "例① 精确入口", "例② 向量入口"]
    rows = [
        ("问题", "IND 要求、NDA 不要求的法规", "细胞治疗相关且属于 IND 的法规"),
        ("进门零件", "Resolve（实体索引精确查）", "Seed（块向量库 ANN）"),
        ("查向量库？", "否", "是（P1）"),
        ("关键文字", "直接用名字 IND / NDA", "“细胞治疗 干细胞 免疫细胞 …”"),
        ("集合运算", "Diff（差集）", "Intersect（交集）"),
        ("共同优化", "惰性取原文 · 并行 · 中间只用 id · 逐条溯源", "同左"),
    ]
    x0, y0 = 70, 182
    cw = [180, 480, 480]
    xs = [x0, x0 + cw[0], x0 + cw[0] + cw[1]]
    rh = 60
    # header row
    for j, c in enumerate(cols):
        rct(s, xs[j], y0, cw[j], 44, NAVY)
        if c:
            txt(s, xs[j] + 16, y0, cw[j] - 20, 44, c, 15, "#FFFFFF", True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    for i, r in enumerate(rows):
        y = y0 + 44 + i * rh
        fill = "#F4F8FD" if i % 2 == 0 else PANEL
        for j in range(3):
            rct(s, xs[j], y, cw[j], rh, fill, line=BORDER, lw=0.75)
        txt(s, xs[0] + 16, y, cw[0] - 20, rh, r[0], 14, NAVY, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        txt(s, xs[1] + 16, y, cw[1] - 24, rh, r[1], 14, INK, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        txt(s, xs[2] + 16, y, cw[2] - 24, rh, r[2], 14, INK, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    footer_native(s)


def s_closing(s, idx):
    bg_native(s); header_native(s, "POSITION", "小结与定位", idx)
    txt(s, 70, 140, 1140, 28, "一句话：向量负责“进门找相关”，实体层关系负责“找全找准”，SQG 让思路看得见，fullname 负责对齐与溯源。", 16, MUTE)
    pts = [
        ("对比裸向量 RAG", "不止取 Top-K 相似块，而是编译出可见的查询步骤、在两层图上按意图走：能多跳、能列全、能溯源、能排查。", BLUE),
        ("对比 OG-RAG（超图+集合覆盖）", "我们用两层图 + 大模型编译的算子 DAG + 优化器/执行器，把“怎么查”显式化、可视化，贴合文档天然层级。", TEAL),
        ("落地节奏", "先小本体+固定算子子集跑通；再逐步加优化器能力与算子。评估：计划正确率/覆盖率/漂移率/溯源率。", ORANGE),
    ]
    y = 186
    for t, d, c in pts:
        card_native(s, 70, y, 1140, 116, c)
        txt(s, 94, y + 20, 1090, 28, t, 18, INK, True)
        txt(s, 94, y + 54, 1090, 50, d, 14.5, MUTE)
        y += 134
    footer_native(s)


def s_effort(s, idx):
    bg_native(s); header_native(s, "EFFORT", "工时估算 · 拆分", idx)
    txt(s, 70, 140, 1140, 28, "按模块拆分的工时估算（单位：小时）。合计约 300 小时，核心在「查询编译与执行」。", 16, MUTE)
    soft = {BLUE: BLUE_SOFT, TEAL: TEAL_SOFT, PURPLE: PURPLE_SOFT,
            GREEN: GREEN_SOFT, ORANGE: ORANGE_SOFT, NAVY: "#E1EAF4"}
    phases = [
        ("A", "数据与索引（离线）", 60, BLUE,
         [("语义切块 + fullname", 12), ("实体抽取与归一", 16),
          ("LLM 建边（关系裁定）", 18), ("块向量化 + 入库", 14)]),
        ("B", "存储与检索底座", 40, TEAL,
         [("AI Search 索引（检索/过滤/重排）", 16),
          ("边表 / 图存储 + 遍历接口", 14), ("实体索引（精确/别名/属性）", 10)]),
        ("C", "查询编译与执行（核心）", 80, PURPLE,
         [("SQG 逻辑算子 + LLM 编译", 20), ("优化器（校验/绑定/下推/并行）", 22),
          ("Search 封装 + Lift/Ground/Traverse", 26), ("集合 / LLM 加工 / 校验控制", 12)]),
        ("D", "集成与可视化", 75, GREEN,
         [("Copilot Studio 集成（接入 RAG 检索）", 40), ("端到端串联（问题→答案+溯源）", 14),
          ("算子级 trace / 可排查", 12), ("溯源与 fullname 回填", 9)]),
        ("E", "评估与调优", 25, ORANGE,
         [("评测集（正确率/覆盖率/溯源率）", 12), ("调参 / 回归 / 误差分析", 13)]),
        ("F", "项目管理与缓冲", 20, NAVY,
         [("设计评审 / 文档", 10), ("联调 / 部署 / 缓冲", 10)]),
    ]
    x0, y0, cw, ch, gx, gy = 70, 178, 366, 176, 21, 16
    for i, (lt, name, hrs, c, subs) in enumerate(phases):
        col, row = i % 3, i // 3
        x, y = x0 + col * (cw + gx), y0 + row * (ch + gy)
        card_native(s, x, y, cw, ch, c)
        circle(s, x + 24, y + 18, 32, soft[c])
        txt(s, x + 24, y + 18, 32, 32, lt, 15, c, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        txt(s, x + 66, y + 18, cw - 150, 32, name, 14.5, INK, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        rrect(s, x + cw - 82, y + 20, 64, 27, c, rad=13)
        txt(s, x + cw - 82, y + 20, 64, 27, f"{hrs}h", 14, "#FFFFFF", True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        hline(s, x + 22, y + 56, x + cw - 22, LINE_C, 1)
        yy = y + 68
        for (tn, th) in subs:
            circle(s, x + 28, yy + 6, 6, c)
            txt(s, x + 42, yy, cw - 150, 20, tn, 12, MUTE, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
            txt(s, x + cw - 100, yy, 76, 20, f"{th}h", 12, INK, True, PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)
            yy += 25
    # ---- total banner + stacked proportion bar ----
    rrect(s, 70, 556, 1140, 100, "#0E2E52", rad=14)
    txt(s, 92, 574, 220, 30, "合计 300 小时", 19, "#FFFFFF", True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    txt(s, 290, 574, 900, 30, "≈ 7.5 人周（按 40h/周）　·　含 Copilot Studio 集成 40h；核心「查询编译与执行」80h",
        13, "#A9C4DE", False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    bx, bw, by = 92, 1096, 614
    sc = bw / 300.0
    cx = bx
    for (lt, name, hrs, c, _) in phases:
        w = hrs * sc
        rrect(s, cx, by, w - 4, 28, c, rad=6)
        txt(s, cx, by, w - 4, 28, f"{lt} {hrs}h", 11.5, "#FFFFFF", True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        cx += w


# ---------------- image slides ----------------
def render(key, svg_str):
    hp = os.path.join(OUT, key + ".html"); pp = os.path.join(OUT, key + ".png")
    with open(hp, "w", encoding="utf-8") as f:
        f.write(html_wrap(svg_str))
    if os.path.exists(pp):
        os.remove(pp)
    subprocess.run([CHROME, "--headless", "--disable-gpu", "--hide-scrollbars", "--no-sandbox",
                    "--force-device-scale-factor=2", "--window-size=1280,720",
                    "--screenshot=" + pp, "file:///" + hp.replace("\\", "/")],
                   capture_output=True, timeout=90)
    return pp


def s_image(s, key, fn, kicker, title, idx, prs):
    pp = render(key, fn())
    if os.path.exists(pp) and os.path.getsize(pp) > 0:
        s.shapes.add_picture(pp, 0, 0, width=prs.slide_width, height=prs.slide_height)
    header_native(s, kicker, title, idx)


def s_full(s, key, fn, prs):
    pp = render(key, fn())
    if os.path.exists(pp) and os.path.getsize(pp) > 0:
        s.shapes.add_picture(pp, 0, 0, width=prs.slide_width, height=prs.slide_height)


DECK = [
    ("full", "01_cover", D.cover, None, None, None,
     "封面：法规检索系统设计。两层图 + 大模型编译可见的算子图(SQG) + 优化器/执行器。可见、可溯源、可排查。"),
    ("img", "02_overview", D.overview, "OVERVIEW", "全景设计图", 2,
     "整体一张图：索引阶段把文档建成两层图并落入三存储；查询阶段大模型编译算子图→优化器→执行器读存储→答案。"),
    ("img", "03_philosophy", D.philosophy, "PRINCIPLES", "核心设计理念", 3,
     "六条理念：可见的算子图、两层图、逻辑与执行分离、关系只用LLM建、优化后执行、可溯源。"),
    ("img", "04_two_layer", D.two_layer, "DATA MODEL", "数据模型 · 两层图", 4,
     "实体层=概念节点(前缀=类型,结构关系挂这)；块层=原文片段(fullname,向量入口)；evidence 边连两层。关系存一次，反向靠遍历。"),
    ("img", "04b_stores", D.stores_accuracy, "STORES", "三个存储 · 为什么能提准确率", 4,
     "三存储各司其职并接力：AI Search(向量/块)语义召回找相关；实体索引把别名/变体归一到唯一节点(消歧、发稳定id)；边表/图沿类型关系确定性遍历(找全可核对)。机制:向量的模糊被实体归一收敛,实体的关系不清被图补上,图的证据在哪被块原文补上。召回×精确×完整×溯源=高正确率。"),
    ("img", "04c_entity_store", D.entity_store, "STORE · 实体索引", "实体索引 · 长什么样 / 为什么这么存 / 怎么检索", 4,
     "实体索引=以名称/别名为键的概念档案表,每个概念只存一条记录(规范名+全部别名+属性+稳定id)。为什么:去重归一、别名归一、属性可筛、发id给别处引用。检索:名字/别名→规范化→哈希精确查O(1)→唯一id;未命中查别名表;按属性枚举/过滤。"),
    ("img", "04d_edge_store", D.edge_store, "STORE · 边表/图", "边表 / 图 · 长什么样 / 为什么这么存 / 怎么检索", 4,
     "边表/图=概念间有向带类型的关系。共五种边(中文名):要求requires(申报→法规)、归类belongs_to(法规→分类)、替代supersedes(新版→旧版)、引用references(文本提到另一条)、出处evidence(概念↔原文块,双向Ground/Lift)。每种在图上颜色标示一个示例。为什么:关系只存一次(有向一行)无冗余、带类型可分别遍历、带权支持扩散排序、反向靠遍历不另存。检索:沿某类型边Traverse多跳;块↔实体走出处边(Lift/Ground)。"),
    ("img", "04e_block_store", D.block_store, "STORE · AI Search", "AI Search · 长什么样 / 为什么这么存 / 怎么检索", 4,
     "AI Search=原文切块+向量+fullname+元数据,一个索引兼做检索/过滤/重排。为什么:小块=可引用最小单元、fullname可溯源可排序、向量语义相邻能召回近义、元数据支持过滤下推。检索:embed查询→HNSW近邻(ANN)+$filter下推早筛→语义重排取Top-k→按fullname溯源或Lift抬到实体。"),
    ("img", "05_index", D.index_flow, "INDEX", "索引阶段", 5,
     "索引阶段四步：切块赋fullname→抽实体归一→LLM建边→块向量化，落入三存储。建边只用LLM(候选提名+LLM裁定)。"),
    ("img", "06_query", D.query_pipeline, "QUERY", "查询阶段 · 三段式", 6,
     "查询三段式：SQG逻辑算子图(大模型)→优化器(校验/绑定/优化)→执行器(物理算子,读三存储)→答案+溯源。检索侧塔缩成一个 Search 算子(=一次 AI Search 调用),精确/语义只是参数不同。"),
    ("img", "07_logical", D.logical_ops, "OPERATORS", "逻辑算子目录", 7,
     "七个逻辑算子：检索/关联/筛选/对比/汇总/校验/回答。意图级、一句人话、固定目录。"),
    ("img", "08_mapping", D.physical_map, "MAPPING", "物理算子 + 逻辑→物理映射", 8,
     "物理算子目录 + 逻辑→物理映射。一个逻辑算子展开成多个物理算子，且因输入不同有不同展开(检索的精确/语义两种)。"),
    ("img", "09_optimizer", D.optimizer, "OPTIMIZER", "优化器做什么", 9,
     "优化器：校验→绑定→优化(下推/并行/惰性取原文/缓存)→编排执行。校验不过可有界回环。"),
    ("img", "10_ex1_plans", D.ex1_plans, "EXAMPLE 1", "例① 精确入口 · SQG + 物理计划", 10,
     "例①“IND要求、NDA不要求的法规”。SQG:检索×2→对比→回答；PEP:Search(精确名)×2→Traverse×2→Diff→Ground→Answer。全程走关键词不查向量,惰性取原文。"),
    ("img", "11_ex1_exec", D.ex1_exec, "EXAMPLE 1", "例① 逐步执行 + 结果", 11,
     "执行:IND集(10)、NDA集(28)只用id→Diff得IND独有8条→最终才Ground取原文→回答附fullname。"),
    ("img", "12_ex2_plans", D.ex2_plans, "EXAMPLE 2", "例② 向量入口 · SQG + 物理计划", 12,
     "例②“细胞治疗相关且属于IND的法规”。op1语义走向量(Search向量→Lift),op2精确走Search(精确名)。同一 Search 算子两种参数同框,在实体层求交集。关键文字:细胞治疗 干细胞…"),
    ("img", "13_ex2_exec", D.ex2_exec, "EXAMPLE 2", "例② 逐步执行 + 结果", 13,
     "执行:P1查向量库命中块→Lift得细胞治疗类Reg(约10)；与IND法规集(10)求交集=7条；Ground取原文→回答。"),
    ("img", "14_compare", D.compare, "SUMMARY", "两个例子对比", 14,
     "两例对比:同一 Search 算子的精确 vs 向量两种参数;是否查向量;关键文字;集合运算Diff vs Intersect;共同优化(惰性/并行/溯源)。"),
    ("img", "15_closing", D.closing, "POSITION", "小结与定位", 15,
     "定位:对比裸RAG(可见/多跳/列全/溯源)、对比OG-RAG(两层图+算子DAG)；落地节奏与评估指标。"),
    ("native", "15b_effort", s_effort, "EFFORT", "工时估算 · 拆分", 15,
     "工时合计约300小时(≈7.5人周)。拆分:A数据与索引60h、B存储与检索底座40h、C查询编译与执行(核心)80h、D集成与可视化75h(含Copilot Studio集成RAG检索40h)、E评估与调优25h、F项目管理与缓全20h。本页为可编辑原生文本。"),
    ("img", "16_algo_overview", D.algo_overview, "ALGORITHMS", "物理算子算法总览", 16,
     "附录：逐个讲解每个物理算子。检索侧塌缩成一个 Search 算子(=一次 AI Search 调用)；只自建 AI Search 不做的四类：图导航、集合、LLM加工、校验控制。"),
    ("img", "17_algo_search", D.algo_search, "ALGORITHM", "Search · 一次 AI Search 调用", 17,
     "Search=一次 POST /docs/search：向量/关键词(search+vectorQueries)、别名(Synonym Map)、过滤($filter)、语义重排(queryType=semantic 重排Top-50)、top/排序/阈值。旧Resolve/Seed/ScanEntities/FilterField/Rerank/TopK/Sort/Threshold全是它的字段,合并成一个算子。"),
    ("img", "18_algo_lift", D.algo_lift_ground, "ALGORITHM", "Lift / Ground · 块↔实体", 18,
     "Lift：块→沿evidence反向→实体(去重)；Ground：实体→沿evidence正向→证据块。数据源：边表+块存储。AI Search 不做,自建。"),
    ("img", "19_algo_traverse", D.algo_traverse, "ALGORITHM", "Traverse · 带权扩散", 19,
     "Traverse=Spreading Activation：act(seed)=1,沿边BFS act(n)=act(cur)×w×γ,act<τ或超跳数停,visited防环。多跳图遍历是AI Search非图库、必须自建的核心。例IND requires→GCP act=1×0.9×0.6=0.54。"),
    ("img", "20_algo_setops", D.algo_setops, "ALGORITHM", "集合运算 · Intersect / Diff / Union", 20,
     "集合运算：哈希集合,Intersect(交)/Diff(差)/Union(并)/Dedup(去重),≈O(|A|+|B|),纯内存跨结果集,AI Search不做。例①Diff、例②Intersect。"),
    ("img", "21_algo_verify", D.algo_verify, "ALGORITHM", "校验控制 · SetCheck / Provenance / Abstain", 21,
     "校验：SetCheck(结果vs应有集合求差找漏项)、ProvenanceCheck(名字在原文字符串匹配,幻觉剔除)、Abstain(分<τ低置信兜底)。不达标可Replan回环。"),
    ("img", "22_algo_llm", D.algo_llm, "ALGORITHM", "LLM 加工算子", 22,
     "LLM算子=固定代码+模板提示词+低温→结构化输出(安全,非任意代码),走Azure OpenAI。Summarize/Extract/Compare/Generate/Judge统一模式；Generate强约束逐条+溯源。"),
]


def build():
    prs = Presentation()
    prs.slide_width = PX(1280); prs.slide_height = PX(720)
    blank = prs.slide_layouts[6]
    for kind, key, fn, kicker, title, idx, notes in DECK:
        s = prs.slides.add_slide(blank)
        if kind == "full":
            s_full(s, key, fn, prs)
        elif kind == "img":
            s_image(s, key, fn, kicker, title, idx, prs)
        else:
            fn(s, idx)
        s.notes_slide.notes_text_frame.text = notes
        print("built", key, kind)
    prs.save(PPTX)
    print("SAVED", os.path.basename(PPTX), len(DECK), "slides")


if __name__ == "__main__":
    build()
