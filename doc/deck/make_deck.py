# -*- coding: utf-8 -*-
"""Render the Assertion-first architecture deck and build an editable-title PPTX."""
from __future__ import annotations

import os
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

import diagrams
from slides_lib import html_wrap

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
OUTPUT = ROOT.parent / "法规检索系统设计.pptx"
CHROME_CANDIDATES = [
    Path(r"C:\Program Files\Google\Chrome\Application\chrome.exe"),
    Path(r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"),
    Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"),
    Path(r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"),
]

NAVY = RGBColor(18, 58, 99)
MUTE = RGBColor(90, 107, 123)
LINE = RGBColor(216, 225, 236)
WHITE = RGBColor(255, 255, 255)
SECTION_COLORS = {
    "系统全景": RGBColor(20, 112, 196),
    "事实与存储": RGBColor(122, 92, 192),
    "索引构建": RGBColor(224, 129, 46),
    "查询执行": RGBColor(14, 156, 156),
    "运行与实现": RGBColor(46, 158, 107),
}

SLIDES = [
    {
        "key": "cover",
        "title": "法规检索系统设计",
        "section": "封面",
        "fn": diagrams.cover,
        "notes": "Assertion-first 最终架构。\n本版完整替换旧的粗粒度关系模型。\n核心主线：Block 证据、Assertion 事实、稳定词汇、派生 Graph、Generation 发布、SQG/PEP 查询。",
    },
    {
        "key": "overview",
        "title": "系统全景",
        "section": "系统全景",
        "fn": diagrams.overview,
        "notes": "先建立端到端心智模型。\n索引负责把法规原文变成可发布事实；查询负责把自然语言编译成可见计划。\nGraph 是导航结构，答案最终仍回到 Assertion 和原文。",
    },
    {
        "key": "principles",
        "title": "核心设计原则",
        "section": "系统全景",
        "fn": diagrams.principles,
        "notes": "六条原则共同解决旧模型的根本问题。\n事实与导航分离、构建与发布分离、意图与执行分离。\nCollection 和 Generation Scope 共同形成可证明的查询边界。",
    },
    {
        "key": "truth_model",
        "title": "从原文到可导航事实",
        "section": "事实与存储",
        "fn": diagrams.truth_model,
        "notes": "Block 保存原文；Assertion 保存法律事实；Entity 和 Action 提供稳定比较键。\nGraph Edge 必须通过 graph_edge_support 回到支持 Assertion。\n任何结论都要能解释事实和原文两层依据。",
    },
    {
        "key": "assertion_anatomy",
        "title": "一条法规事实如何表达",
        "section": "事实与存储",
        "fn": diagrams.assertion_anatomy,
        "notes": "Assertion 把主体、行动、模态、条件、例外和证据放在一个原子事实中。\n这避免把条件许可压缩成模糊关系。\nAction 本身不包含模态，便于跨主体做集合比较。",
    },
    {
        "key": "schema_map",
        "title": "数据存储与表结构",
        "section": "事实与存储",
        "fn": diagrams.schema_map,
        "notes": "Azure SQL 保存结构事实、发布状态和运行记录。\nAzure AI Search 保存 Block 原文与向量。\n两侧都带 generation_id，查询只读取冻结代次。",
    },
    {
        "key": "generation",
        "title": "Generation 原子发布",
        "section": "事实与存储",
        "fn": diagrams.generation,
        "notes": "索引构建采用完整隔离代次。\n质量通过之前，新数据不可查；失败或取消也不污染线上。\n激活事务一次更新新旧代次状态和 Store 指针。",
    },
    {
        "key": "index_workflow",
        "title": "索引工作流",
        "section": "索引构建",
        "fn": diagrams.index_workflow,
        "notes": "解析后，向量写入和逐块抽取并行。\n归一、持久化和 Graph 派生集中执行。\n通用 Workflow 负责依赖、取消、状态和 Token 记录。",
    },
    {
        "key": "extraction_guardrails",
        "title": "抽取校验与隔离",
        "section": "索引构建",
        "fn": diagrams.extraction_guardrails,
        "notes": "LLM 输出不直接入库。\n先做反馈重试和可确定的本地修复，再逐条隔离坏项。\n单条错误不拖垮全文，但不可定位 Quote 等硬错误仍会阻止发布。",
    },
    {
        "key": "normalization",
        "title": "Entity / Action 归一",
        "section": "索引构建",
        "fn": diagrams.normalization,
        "notes": "Entity 先按同类型规范名精确匹配，别名只作回退。\nAction 使用包含对象、接收方和限定语的完整签名。\n目标是高精度和可解释，而不是最大化合并。",
    },
    {
        "key": "quality_gate",
        "title": "质量门禁",
        "section": "索引构建",
        "fn": diagrams.quality_gate,
        "notes": "质量门禁是激活事务的前置条件。\n硬错误必须为零；隔离比例受控时可带警告发布。\n门禁失败时，Store 继续指向旧 Active Generation。",
    },
    {
        "key": "collection_scope",
        "title": "Collection 与查询快照",
        "section": "查询执行",
        "fn": diagrams.collection_scope,
        "notes": "Initializer 在任何读取前选择一个可见 Collection。\n随后冻结成员 Store 与各自 Active Generation。\nSQL、Search、Graph 和文档绑定都必须复用这个不可变快照。",
    },
    {
        "key": "query_stages",
        "title": "查询五阶段",
        "section": "查询执行",
        "fn": diagrams.query_stages,
        "notes": "五个阶段职责固定且各自持久化。\nLLM 只用于 SQG 理解和最终表达；PEP 与算子执行保持确定性。\n生成器没有 Repository，因此不能在回答时再次检索。",
    },
    {
        "key": "sqg_intents",
        "title": "强类型 SQG",
        "section": "查询执行",
        "fn": diagrams.sqg_intents,
        "notes": "SQG 只表达业务意图，不出现 Store、Generation、TopK 或物理算子。\n所有命名对象必须唯一绑定当前可见 Catalog。\n首次校验失败反馈重试一次，最终失败保留 Raw Invalid Output。",
    },
    {
        "key": "pep_templates",
        "title": "确定性 PEP",
        "section": "查询执行",
        "fn": diagrams.pep_templates,
        "notes": "Planner 不调用 LLM。\n它先绑定稳定 ID，再按 SQG Intent 选择固定模板。\n节点端口与输出强类型，顶层必须显式绑定 facts 和 evidence。",
    },
    {
        "key": "operators",
        "title": "物理算子",
        "section": "查询执行",
        "fn": diagrams.operators,
        "notes": "算子分为绑定、事实读取、图导航、集合、取证和原文检索。\n集合算子按稳定 Comparison Key 计算并合并来源。\nGroundAssertions 只取结果事实自己的证据。",
    },
    {
        "key": "graph_role",
        "title": "Graph 的职责边界",
        "section": "查询执行",
        "fn": diagrams.graph_role,
        "notes": "Graph 适合关系导航、反查和集合计算。\n每条边都保留 Assertion 支持，可继续下钻到 Quote。\n图中没有一条边，只说明当前冻结快照没有该派生关系。",
    },
    {
        "key": "example_diff",
        "title": "示例：主体职责差集",
        "section": "查询执行",
        "fn": diagrams.example_diff,
        "notes": "两个主体分别执行 EntityLookup 和 SubjectActions。\nDiff 做 A 减 B，再 GroundAssertions 只取幸存事实原文。\n答案必须说明差集只针对当前 Collection 快照。",
    },
    {
        "key": "example_documents",
        "title": "示例：文档比较",
        "section": "查询执行",
        "fn": diagrams.example_documents,
        "notes": "明确文档比较采用每文档独立 BlockSearch。\nEvidenceBundle 保持分组，避免把文档 A 的证据当成文档 B 的依据。\n任一分支没有证据时，生成器明确报告证据不足。",
    },
    {
        "key": "observability",
        "title": "运行记录与失败定位",
        "section": "运行与实现",
        "fn": diagrams.observability,
        "notes": "索引和查询都记录完整 DAG 与阶段状态。\n保存输入输出、Token、耗时、错误和原始非法输出。\n目标是把答案问题直接定位到阶段、算子、事实或 Block。",
    },
    {
        "key": "module_boundaries",
        "title": "模块边界",
        "section": "运行与实现",
        "fn": diagrams.module_boundaries,
        "notes": "最终 Nexus 只有 domain、infrastructure、indexing、querying 四个包。\n通用 Workflow 独立，不知道 Nexus 业务语义。\n旧两套模型已删除，避免运行路径和事实口径分裂。",
    },
    {
        "key": "closing",
        "title": "系统定位",
        "section": "运行与实现",
        "fn": diagrams.closing,
        "notes": "最后回到系统定位。\nBlock 提供证据，Assertion 提供事实，稳定词汇提供比较键，Graph 提供导航。\nGeneration、Collection、SQG 和 PEP 分别保障发布、范围、意图与执行。",
    },
]


def find_browser() -> Path:
    for path in CHROME_CANDIDATES:
        if path.exists():
            return path
    raise FileNotFoundError("Chrome or Edge was not found in the standard installation paths.")


def render_assets() -> list[Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    browser = find_browser()
    rendered: list[Path] = []
    for index, spec in enumerate(SLIDES, start=1):
        html_path = ASSETS / f"slide_{index:02d}.html"
        png_path = ASSETS / f"slide_{index:02d}.png"
        html_path.write_text(html_wrap(spec["fn"]()), encoding="utf-8")
        url = html_path.resolve().as_uri()
        command = [
            str(browser),
            "--headless",
            "--disable-gpu",
            "--hide-scrollbars",
            "--no-sandbox",
            "--force-device-scale-factor=2",
            "--window-size=1280,720",
            f"--screenshot={png_path.resolve()}",
            url,
        ]
        result = subprocess.run(command, capture_output=True, text=True, encoding="utf-8", errors="replace")
        if result.returncode != 0 or not png_path.exists():
            raise RuntimeError(f"Failed to render slide {index}: {result.stderr}")
        rendered.append(png_path)
        print(f"rendered {index:02d}/{len(SLIDES):02d}")

    expected = {f"slide_{i:02d}.html" for i in range(1, len(SLIDES) + 1)} | {
        f"slide_{i:02d}.png" for i in range(1, len(SLIDES) + 1)
    }
    for path in ASSETS.glob("slide_*.*"):
        if path.name not in expected:
            path.unlink(missing_ok=True)
    return rendered


def add_text(slide, x, y, w, h, value, size, color, bold=False, font="Microsoft YaHei",
             align=PP_ALIGN.LEFT, margin=0, valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def build_pptx(images: list[Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "法规检索系统设计 · Assertion-first 最终架构"
    prs.core_properties.subject = "Block 证据、Legal Assertion、稳定词汇、派生 Graph、Generation 发布、强类型 SQG 与确定性 PEP"
    prs.core_properties.author = "NexusRAG"
    prs.core_properties.keywords = "法规检索, Assertion-first, Generation, SQG, PEP, Graph"

    for index, (spec, image) in enumerate(zip(SLIDES, images), start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        picture = slide.shapes.add_picture(str(image), 0, 0, width=prs.slide_width, height=prs.slide_height)
        picture.name = f"Architecture diagram · {spec['key']}"

        if index > 1:
            color = SECTION_COLORS[spec["section"]]
            add_text(slide, 0.73, 0.20, 4.8, 0.20, spec["section"].upper(), 8.5, color,
                     bold=True, font="Segoe UI", valign=MSO_ANCHOR.MIDDLE)
            add_text(slide, 0.73, 0.42, 11.2, 0.60, spec["title"], 25.5, NAVY,
                     bold=True, valign=MSO_ANCHOR.MIDDLE)
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.73), Inches(1.20), Inches(0.62), Inches(0.045))
            accent.fill.solid()
            accent.fill.fore_color.rgb = color
            accent.line.fill.background()
            rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.37), Inches(1.20), Inches(11.22), Inches(0.018))
            rule.fill.solid()
            rule.fill.fore_color.rgb = LINE
            rule.line.fill.background()
            add_text(slide, 12.22, 0.32, 0.36, 0.26, f"{index:02d}", 9, MUTE,
                     bold=True, font="Segoe UI", align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)

        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.text = spec["notes"]

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"saved {OUTPUT.name} ({len(SLIDES)} slides)")


def main() -> int:
    try:
        images = render_assets()
        build_pptx(images)
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
